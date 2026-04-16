import os
from pptx import Presentation
from fpdf import FPDF
import re

def extract_syllabus():
    ppt_dir = '/Users/jyothssena/CourseLens/CourseLens_data/ppts/'
    files = sorted([f for f in os.listdir(ppt_dir) if f.endswith('.pptx')])
    
    # Sort files naturally (e.g. chap01, chap02)
    def extract_num(f):
        m = re.search(r'\d+', f)
        return int(m.group()) if m else 0
    
    files = sorted(files, key=extract_num)

    syllabus = []
    chapter_no=1
    for file in files:
        filepath = os.path.join(ppt_dir, file)
        try:
            prs = Presentation(filepath)
            chapter_title = f"{file}"
            topics = []
            
            for i, slide in enumerate(prs.slides):
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        if i == 0:
                            chapter_title = title_text
                        else:
                            if title_text not in topics:
                                topics.append(title_text)
            
            syllabus.append({
                "chapter": chapter_title,
                "number":chapter_no,
                "topics": topics
            })
            chapter_no+=1
        except Exception as e:
            print(f"Error reading {file}: {e}")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    # Add a Unicode font
    # fpdf2 natively supports utf-8 but needs a font that supports it if there are special chars
    # We will just use Helvetica and replace problematic chars
    pdf.set_font("Helvetica", style="B", size=16)
    pdf.cell(0, 10, "Course Syllabus", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(10)
    
    for item in syllabus:
        pdf.set_font("Helvetica", style="B", size=14)
        ch_title = item['chapter'].encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(0, 10, f"Chapter {item['number']}: {ch_title}", new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_font("Helvetica", size=12)
        for topic in item['topics']:
            top_title = topic.encode('latin-1', 'replace').decode('latin-1')
            # remove newlines
            top_title = top_title.replace('\n', ' ').replace('\r', '')
            pdf.cell(0, 8, f"    - {top_title}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)
        
    output_path = '/Users/jyothssena/CourseLens/Syllabus.pdf'
    pdf.output(output_path)
    print(f"Syllabus PDF generated at: {output_path}")

if __name__ == "__main__":
    extract_syllabus()
