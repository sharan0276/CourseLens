# Processes the PPTX files and extracts the text content from them.

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from domain.slide import Slide
from services.ingestion.content_extractor import ContentExtractor
from services.ingestion.image_extractor import ImageExtractor
import re


class PPTXParser:
    """
    Resposible for reading PPTX files and extracting text content from them.
    """
    def __init__(self, image_folder_path: str):
        self.content_extractor = ContentExtractor()
        self.image_extractor = ImageExtractor(image_folder_path)

    def parse(self, filepath: str):
        """
        Parses the PPTX file and returns the slides as a list of Slide objects.
        """
        presentation = Presentation(filepath)
        slides = []
        
        # Extracting digits from filepath
        lec_num = int(re.search(r'\d+', filepath).group())
        only_images = False
        for slide_number, slide in enumerate(presentation.slides, start = 1):

            #Extract title
            title = slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {slide_number}"

            txt_blocks = self.content_extractor.extract(slide, slide_number)
            images = self.image_extractor.extract(slide, slide_number, filepath)

            if not txt_blocks :
                only_images = True

            slides.append(Slide(slide_number = slide_number, title = title, content = txt_blocks, images = images, source_file = os.path.basename(filepath), lecture_number = lec_num, source_type = "pptx", only_images = only_images))
            only_images = False                 

        return slides