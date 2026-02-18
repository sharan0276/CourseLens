from pptx.enum.shapes import MSO_SHAPE_TYPE

class ContentExtractor:
    """
    Responsible for extracting text content from PPTX files.
    """
    def __init__(self):
        pass

    def extract(self, slide, slide_number: int) -> list:
        """
        Main Entry point. Seperates Shapes by type and applies precedence to decide which shape to extract text from depending on 
        teh type available on the slide.

        Precedence:
        1. AUTO_SHAPES with content in bullet points
        2. AUTO_SHAPES with multiple paragraphs
        3. Non Title PLACEHOLDERS
        4. GROUPED SHAPES - Unpacked recursively
        """
        
        # Tracker to avoid duplicates
        seen_texts =  set() 

        # Categorize shapes by type
        auto_shapes, placeholders, groups = [], [], []
        for shape in slide.shapes:
            st = shape.shape_type

            if st == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT or st == MSO_SHAPE_TYPE.PICTURE:
                continue
            elif st == MSO_SHAPE_TYPE.GROUP:
                groups.append(shape)
            elif st == MSO_SHAPE_TYPE.PLACEHOLDER:
                placeholders.append(shape)
            # Any other shapes are considered as auto shapes
            elif shape.has_text_frame: 
                auto_shapes.append(shape)

        # Priority 1 : AUTO_SHAPES with level > 0 (bullet points) 
        levelled = [s for s in auto_shapes if self._has_levels(s)]
        if levelled:
            blocks = self._extract_from_shapes(levelled, seen_texts)
            if blocks:
                return blocks

        # Priority 2 : AUTO_SHAPES with multiple paragraphs
        multi_para = [s for s in auto_shapes if self._has_multiple_paragraphs(s)]
        if multi_para:
            blocks = self._extract_from_shapes(multi_para, seen_texts)
            if blocks:
                return blocks

        # Priority 3 : Non Title PLACEHOLDERS
        non_title_placeholders = [p for p in placeholders if not self._is_title_placeholder(p)]
        if non_title_placeholders:
            blocks = self._extract_from_shapes(non_title_placeholders, seen_texts)
            if blocks:
                return blocks

        # Priority 4 : GROUPED SHAPES - Unpacked recursively
        if groups:
            blocks = self._extract_from_groups(groups, seen_texts)
            if blocks:
                return blocks

        
        return []


    def _extract_from_shapes(self, shapes: list, seen_texts: set) -> list:
        """
        Iterates through the grouped shapes to build a parent-child block especially for paragarphs with bullet points.
        To retain heirarchy of bullet points.
        """
        blocks = []

        for shape in shapes:
            if not shape.has_text_frame:
                continue

            current_parent = None

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()

                # Skip empty para and already seen text
                if not text or (text in seen_texts):
                    continue

                total_paras = len([p for p in shape.text_frame.paragraphs if p.text.strip()])
                if len(text) < 15 and total_paras == 1:
                    continue

                seen_texts.add(text)
                level = para.level

                if level == 0:
                    # Start a new parent block
                    # Any subsequent blocks are added as children
                    current_parent = {
                        "level": 0,
                        "text": text,
                        "children": []
                    }
                    blocks.append(current_parent)
                elif current_parent:
                    # Add as child to the current parent
                    current_parent["children"].append({
                        "level": level,
                        "text": text
                    })
                else:
                    # No active parent
                    blocks.append({
                        "level": level,
                        "text": text,
                        "children": []
                    })
        return blocks

    
    def _extract_from_groups(self, groups: list, seen_texts: set) -> list:
        """
        Unpack GROUP shapes recursively to get all inner shapes and extract text from them.
        """

        all_inner_shapes = []
        for group in groups:
            all_inner_shapes.extend(self._unpack_group(group))

        return self._extract_from_shapes(all_inner_shapes, seen_texts)


    def _unpack_group(self, group):
        """
        Recursively unpacks shapes from a group.
        """
        inner_shapes = []
        for shape in group.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                inner_shapes.extend(self._unpack_group(shape))
            else:
                inner_shapes.append(shape)
        return inner_shapes
            

#-------------------------------------------
#         Adding Helper Functions
#-------------------------------------------

    def _has_levels(self, shape) -> bool:
        """
        Returns True if any non-empty paragraph in the shape has levels.
        """
        if not shape.has_text_frame:
            return False
        
        return any(p.level > 0 for p in shape.text_frame.paragraphs if p.text.strip())

    def _has_multiple_paragraphs(self, shape) -> bool:
        """
        Returns True if shape has more than one non-empty paragraph.
        """
        if not shape.has_text_frame:
            return False
        
        return len([p for p in shape.text_frame.paragraphs if p.text.strip()]) > 1


    def _is_title_placeholder(self, shape) -> bool:
        """
        Returns True if shape is a title placeholder.
        """
        if not shape.has_text_frame:
            return False
        
        return shape.placeholder_format.idx == 0
