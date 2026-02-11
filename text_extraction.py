from PyPDF2 import PdfReader
from pptx import Presentation
import os

def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    
    # Get number of pages
    num_pages = len(reader.pages)
    print(f"Total pages: {num_pages}")
    
    # Extract text from all pages
    full_text = ""
    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()
        full_text += text
        print(f"\n--- Page {page_num + 1} ---")
        print(text)
    
    return full_text

def extract_text(file_path):
    """Automatically detects file type and extracts text"""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)

def save_text_to_file(input_path, output_path):
    _, ext = os.path.splitext(input_path)
    ext = ext.lower()
    
    with open(output_path, 'w', encoding='utf-8') as f:
        if ext == '.pdf':
            reader = PdfReader(input_path)
            for page in reader.pages:
                text = page.extract_text()
                f.write(text)
                f.write('\n\n--- Page Break ---\n\n')
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(input_path)
            for slide_num, slide in enumerate(prs.slides):
                f.write(f"--- Slide {slide_num + 1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
                f.write('\n\n')
    
    print(f"Text saved to {output_path}")

# Usage examples
if __name__ == "__main__":
    # Works with both PDFs and PPTs
    file_paths = ["bw01", "bw02", "bw03", "bw04-conditional", "bw04-iterative", "bw05","bw06",
                  "bw07", "bw08", "bw09", "bw10", "bw11", "bw13", "bw14"]
    
    for i in file_paths:
        save_text_to_file(f'CourseLens_data/pdfs/{i}.pdf', f"CourseLens_data/text_files/{i}")